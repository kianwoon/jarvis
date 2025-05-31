"""
PowerPoint document handler for extracting content from PPT/PPTX files
"""
import os
import re
from typing import List, Optional, Dict, Any, Tuple
from datetime import datetime
import hashlib

from app.document_handlers.base import DocumentHandler, ExtractedChunk, ExtractionPreview
from utils.deduplication import hash_text


class PowerPointHandler(DocumentHandler):
    """Handler for Microsoft PowerPoint presentations (PPT/PPTX)"""
    
    SUPPORTED_EXTENSIONS = ['.ppt', '.pptx']
    
    def __init__(self):
        super().__init__()
        try:
            from pptx import Presentation
            self.Presentation = Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint document processing. "
                "Install it with: pip install python-pptx"
            )
    
    def validate(self, file_path: str) -> Tuple[bool, Optional[str]]:
        """Validate if PowerPoint document can be processed"""
        if not os.path.exists(file_path):
            return False, "File not found"
            
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > self.MAX_FILE_SIZE_MB:
            return False, f"File too large: {file_size_mb:.1f}MB (max: {self.MAX_FILE_SIZE_MB}MB)"
            
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.SUPPORTED_EXTENSIONS:
            return False, f"Unsupported file extension: {file_ext}"
            
        return True, None
    
    def extract(self, file_path: str, options: Dict[str, Any] = None) -> List[ExtractedChunk]:
        """Extract text chunks from PowerPoint presentation"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file extension: {file_ext}")
        
        # For .ppt files, we might need to convert to .pptx first
        if file_ext == '.ppt':
            raise NotImplementedError(
                "Direct .ppt file support not implemented. "
                "Please convert to .pptx format or use LibreOffice conversion."
            )
        
        # Extract text from PPTX
        try:
            pres = self.Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Failed to load PowerPoint presentation: {str(e)}")
        
        file_name = os.path.basename(file_path)
        with open(file_path, "rb") as f:
            file_content = f.read()
            file_id = hashlib.sha256(file_content).hexdigest()[:12]
        
        chunks = []
        
        # Group slides into chunks (e.g., 3-5 slides per chunk for context)
        if options is None:
            options = {}
        slides_per_chunk = options.get('slides_per_chunk', 3)
        include_notes = options.get('include_notes', True)
        
        slide_groups = []
        current_group = []
        
        for slide_idx, slide in enumerate(pres.slides):
            slide_content = self._extract_slide_content(slide, slide_idx + 1, include_notes)
            
            if slide_content['text'] or slide_content['notes']:
                current_group.append(slide_content)
                
                # Check if we should create a new group
                if len(current_group) >= slides_per_chunk:
                    slide_groups.append(current_group)
                    current_group = []
        
        # Add remaining slides
        if current_group:
            slide_groups.append(current_group)
        
        # Convert slide groups to chunks
        for group_idx, slide_group in enumerate(slide_groups):
            chunk_content = self._format_slide_group(slide_group)
            
            if not chunk_content.strip():
                continue
            
            # Calculate quality score
            quality_score = self.calculate_quality_score(chunk_content)
            if quality_score < self.MIN_QUALITY_SCORE:
                continue
            
            # Determine slide range
            first_slide = slide_group[0]['slide_number']
            last_slide = slide_group[-1]['slide_number']
            slide_range = f"{first_slide}-{last_slide}" if first_slide != last_slide else str(first_slide)
            
            # For PowerPoint, slide numbers can directly translate to pages
            # Use the first slide number of the group as the page number
            page_number = slide_group[0]['slide_number']
            
            metadata = {
                'source': file_name.lower(),
                'doc_type': 'powerpoint',
                'section': f"Slides {slide_range}",
                'chunk_index': group_idx,
                'uploaded_at': datetime.now().isoformat(),
                'file_id': file_id,
                'chunk_type': 'slides',
                'slide_range': slide_range,
                'slide_count': len(slide_group),
                'quality_score': quality_score,
                'page': page_number
            }
            
            # Add title if available
            if slide_group[0].get('title'):
                metadata['slide_title'] = slide_group[0]['title']
            
            metadata['doc_id'] = f"{file_id}_slides_{slide_range.replace('-', '_')}"
            metadata['hash'] = hash_text(chunk_content)
            
            chunk = ExtractedChunk(
                content=chunk_content,
                metadata=metadata
            )
            chunks.append(chunk)
        
        return chunks
    
    def _extract_slide_content(self, slide, slide_number: int, include_notes: bool) -> Dict[str, Any]:
        """Extract all text content from a slide"""
        slide_data = {
            'slide_number': slide_number,
            'title': '',
            'text': [],
            'notes': '',
            'tables': []
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()
        
        # Extract text from all shapes
        for shape in slide.shapes:
            # Skip if it's the title (already extracted)
            if shape == slide.shapes.title:
                continue
            
            # Extract text from text frames
            if shape.has_text_frame:
                text_content = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_content.append(para_text)
                
                if text_content:
                    slide_data['text'].extend(text_content)
            
            # Extract text from tables
            if shape.has_table:
                table_text = self._extract_table_from_shape(shape.table)
                if table_text:
                    slide_data['tables'].append(table_text)
        
        # Extract speaker notes if requested
        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data['notes'] = notes_text
        
        return slide_data
    
    def _extract_table_from_shape(self, table) -> str:
        """Extract text from a PowerPoint table"""
        rows = []
        
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_text.append(cell_text)
            
            if any(cell for cell in row_text):
                rows.append(" | ".join(row_text))
        
        if rows:
            return "\n".join(rows)
        return ""
    
    def _format_slide_group(self, slides: List[Dict[str, Any]]) -> str:
        """Format a group of slides into readable text"""
        formatted_parts = []
        
        for slide in slides:
            slide_parts = []
            
            # Add slide header
            slide_header = f"=== Slide {slide['slide_number']}"
            if slide['title']:
                slide_header += f": {slide['title']}"
            slide_header += " ==="
            slide_parts.append(slide_header)
            
            # Add slide content
            if slide['text']:
                slide_parts.append("\n".join(slide['text']))
            
            # Add tables
            for table in slide['tables']:
                slide_parts.append("\n[Table]")
                slide_parts.append(table)
            
            # Add notes
            if slide['notes']:
                slide_parts.append("\n[Speaker Notes]")
                slide_parts.append(slide['notes'])
            
            if len(slide_parts) > 1:  # More than just the header
                formatted_parts.append("\n\n".join(slide_parts))
        
        return "\n\n\n".join(formatted_parts)
    
    def get_preview(self, file_path: str, max_chunks: int = 5) -> ExtractionPreview:
        """Get a preview of extraction without full processing"""
        try:
            # Extract a few chunks for preview
            chunks = self.extract(file_path, options={'slides_per_chunk': 2})[:max_chunks]
            
            # Get file metadata
            file_size = os.path.getsize(file_path)
            preview_data = self.extract_preview(file_path)
            
            metadata = {
                'filename': os.path.basename(file_path).lower(),
                'file_size_mb': round(file_size / (1024 * 1024), 2),
                'doc_type': 'powerpoint',
                **preview_data.get('metadata', {})
            }
            
            return ExtractionPreview(
                total_chunks=len(chunks),
                sample_chunks=chunks,
                file_metadata=metadata
            )
            
        except Exception as e:
            return ExtractionPreview(
                total_chunks=0,
                sample_chunks=[],
                file_metadata={'error': str(e)},
                warnings=[f"Preview generation failed: {str(e)}"]
            )
    
    def extract_preview(self, file_path: str, max_slides: int = 3) -> Dict[str, Any]:
        """Extract preview information from PowerPoint presentation"""
        try:
            pres = self.Presentation(file_path)
            
            preview_slides = []
            total_slides = len(pres.slides)
            
            # Get first few slides for preview
            for slide_idx, slide in enumerate(pres.slides[:max_slides]):
                slide_content = self._extract_slide_content(slide, slide_idx + 1, include_notes=False)
                
                preview_slide = {
                    'slide_number': slide_content['slide_number'],
                    'title': slide_content['title'],
                    'content': '\n'.join(slide_content['text'][:3])  # First 3 text elements
                }
                
                if slide_content['tables']:
                    preview_slide['has_tables'] = True
                
                preview_slides.append(preview_slide)
            
            # Get presentation properties
            properties = {}
            try:
                core_props = pres.core_properties
                if core_props.author:
                    properties['author'] = core_props.author
                if core_props.title:
                    properties['title'] = core_props.title
                if core_props.subject:
                    properties['subject'] = core_props.subject
                if core_props.created:
                    properties['created'] = core_props.created.isoformat()
                if core_props.modified:
                    properties['modified'] = core_props.modified.isoformat()
            except:
                pass
            
            # Count various elements
            text_shapes = 0
            tables = 0
            images = 0
            
            for slide in pres.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_shapes += 1
                    if shape.has_table:
                        tables += 1
                    if hasattr(shape, 'image'):
                        images += 1
            
            return {
                'preview_slides': preview_slides,
                'metadata': {
                    'total_slides': total_slides,
                    'text_shapes': text_shapes,
                    'tables': tables,
                    'images': images,
                    'properties': properties
                }
            }
            
        except Exception as e:
            return {
                'preview_slides': [],
                'metadata': {
                    'error': f"Error loading preview: {str(e)}"
                }
            }